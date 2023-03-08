from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os


class CriadorDeSlides:
    dir_at = os.path.dirname(os.path.realpath(__file__))

    prs = Presentation(dir_at + '\\Padrao.pptx')  # Tema do Slide
    
    FontNameSM = "K2D SemiBold"
    FontNameST = "K2D ExtraBold"

    def __init__(self):
        pass

    def CriaTitulo(self, titulo, sub):
        slideTitulo = CriadorDeSlides.prs.slides.add_slide(CriadorDeSlides.prs.slide_layouts[6])
        x = Cm(20)
        y = Cm(20)
        largura = Cm(50)
        altura = Cm(10)
        cx_titulo = slideTitulo.shapes.add_textbox(x, y, largura, altura)
        cx_titulo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Alinhar a caixa do texto
        t = cx_titulo.text_frame.paragraphs[0]  # t = objeto Slide
        run = t.add_run()
        run.text = titulo  # Titulo
        cx_titulo.text_frame.paragraphs[0].aligment = PP_ALIGN.CENTER
        # AdicionarFont
        font = run.font
        font.name = CriadorDeSlides.FontNameST
        font.size = Pt(250)
        font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # SUBTITULO
        Subx = Cm(20)
        Suby = Cm(40)
        Sublargura = Cm(50)
        Subaltura = Cm(5)
        Subcx_titulo = slideTitulo.shapes.add_textbox(Subx, Suby, Sublargura, Subaltura)
        Subcx_titulo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Alinhar a caixa do texto
        s = Subcx_titulo.text_frame.paragraphs[0]  # t = objeto Slide
        Subrun = s.add_run()
        Subrun.text = sub  # Subitulo
        Subcx_titulo.text_frame.paragraphs[0].aligment = PP_ALIGN.CENTER
        # AdicionarFont
        Subfont = Subrun.font
        Subfont.name = CriadorDeSlides.FontNameST
        Subfont.size = Pt(105)
        Subfont.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def CriaSlideLetra(self, lista, titulo):
        x1 = Cm(10)
        y1 = Cm(13.5)
        largura2 = Cm(70)
        altura2 = Cm(30)
        TamanhoLetra = len(lista)
        # print(TamanhoLetra)
        for linhas in range(TamanhoLetra):
            # print(linhas) 0 1 2 3
            slideLetra = CriadorDeSlides.prs.slides.add_slide(CriadorDeSlides.prs.slide_layouts[6])
            cx_Letra = slideLetra.shapes.add_textbox(x1, y1, largura2, altura2)
            text_frame = cx_Letra.text_frame

            cx_Letra.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            f = cx_Letra.text_frame.paragraphs[0]
            run = f.add_run()

            # paragrafo = text_frame.add_paragraph()
            run.text = lista[linhas]
            cx_Letra.text_frame.paragraphs[0].aligmnet = PP_ALIGN

            font = run.font
            font.name = CriadorDeSlides.FontNameSM
            font.size = Pt(150)
            font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            CriadorDeSlides.prs.save(titulo + ".pptx")
